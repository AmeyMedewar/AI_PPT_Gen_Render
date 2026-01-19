
import os
import re
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def add_markdown_text(paragraph, text):
    """
    Add text with markdown-style bold (**bold**) support to a pptx paragraph.
    """
    paragraph.clear()  # remove default run
    parts = re.split(r"(\*\*.*?\*\*)", text)

    for part in parts:
        if not part:
            continue
        run = paragraph.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]  # strip **
            run.font.bold = True
        else:
            run.text = part
            run.font.bold = False


class SlideGeneratorGold:
    def __init__(self, output="Output/"):
        self.output_dir = output
        os.makedirs(self.output_dir, exist_ok=True)

    def parse_model_output(self, text):
        slides = []
        current_slide = None

        for line in text.splitlines():
            line = line.strip()
            if not line:
                continue

            if line.lower().startswith("slide") and ":" in line:
                if current_slide:
                    slides.append(current_slide)
                current_slide = {"title": "", "subtitle": "", "author": "", "date": "", "tagline": "", "bullets": []}
                slide_title = line.split(":", 1)[1].strip()
                current_slide["title"] = slide_title
                continue

            elif line.lower().startswith("- title:"):
                current_slide["title"] = line.split(":", 1)[1].strip()
            elif line.lower().startswith("- subtitle:"):
                current_slide["subtitle"] = line.split(":", 1)[1].strip()
            elif line.lower().startswith("- author:"):
                current_slide["author"] = line.split(":", 1)[1].strip()
            elif line.lower().startswith("- date:"):
                current_slide["date"] = line.split(":", 1)[1].strip()
            elif line.lower().startswith("- tagline:"):
                current_slide["tagline"] = line.split(":", 1)[1].strip()
            elif line.startswith("- "):
                bullet = line[2:].strip()
                current_slide["bullets"].append(bullet)

        if current_slide:
            slides.append(current_slide)

        return slides


    def create_slide_deck(self, slides_data, filename="generated_ppt.pptx"):
        pres = Presentation()
        filename = slides_data[0].get("title"+".pptx","generated_ppt.pptx")

        for idx, slide_info in enumerate(slides_data):
            layout = pres.slide_layouts[0] if idx == 0 else pres.slide_layouts[1]
            slide = pres.slides.add_slide(layout)

            # Title
            title_shape = slide.shapes.title
            title_para = title_shape.text_frame.paragraphs[0]
            add_markdown_text(title_para, slide_info.get("title", "Untitled Slide"))
            title_para.font.size = Pt(40)
            title_para.font.color.rgb = RGBColor(48, 40, 36)

            # First slide: subtitle, author, date, tagline
            if idx == 0:
                subtitle_placeholder = slide.placeholders[1]
                subtitle_placeholder.text = ""
                subtitle = slide_info.get("subtitle", "")
                if subtitle:
                    p = subtitle_placeholder.text_frame.paragraphs[0]
                    add_markdown_text(p, subtitle)
                    p.font.size = Pt(28)
                    p.font.color.rgb = RGBColor(80, 80, 80)

                # Author + Date
                author_text = f"Author: {slide_info.get('author', 'Amey G M')}\nDate: {slide_info.get('date', 'October 26, 2023')}"
                left_box = slide.shapes.add_textbox(left=Pt(20), top=Pt(400), width=Pt(300), height=Pt(50))
                tf = left_box.text_frame
                add_markdown_text(tf.paragraphs[0], author_text)
                for p in tf.paragraphs:
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(0, 51, 102)

                # Tagline at bottom center
                tagline = slide_info.get("tagline", "")
                if tagline:
                    tagline_box = slide.shapes.add_textbox(left=Pt(100), top=Pt(460), width=Pt(800), height=Pt(50))
                    tf_tag = tagline_box.text_frame
                    add_markdown_text(tf_tag.paragraphs[0], tagline)
                    tf_tag.paragraphs[0].font.size = Pt(18)
                    tf_tag.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
                    tf_tag.paragraphs[0].alignment = PP_ALIGN.CENTER

            # Content slides with bullets
            else:
                content_placeholder = None
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.idx != 0 and hasattr(placeholder, "text_frame"):
                        content_placeholder = placeholder
                        break

                if content_placeholder and content_placeholder.text_frame:
                    tf = content_placeholder.text_frame
                    tf.clear()

                    for bullet in slide_info.get("bullets", []):
                        if ":" in bullet:
                            heading, content = bullet.split(":", 1)
                            p = tf.add_paragraph()
                            add_markdown_text(p, heading.strip() + ":")
                            p.font.size = Pt(20)
                            p.font.color.rgb = RGBColor(0, 0, 0)

                            p_sub = tf.add_paragraph()
                            add_markdown_text(p_sub, content.strip())
                            p_sub.font.size = Pt(18)
                            p_sub.level = 1
                        else:
                            p = tf.add_paragraph()
                            add_markdown_text(p, bullet.strip())
                            p.font.size = Pt(20)
                            p.level = 0

        # Add final Thank You slide at the end
        blank_slide_layout = pres.slide_layouts[6]
        thank_you_slide = pres.slides.add_slide(blank_slide_layout)

        textbox = thank_you_slide.shapes.add_textbox(left=Inches(1), top=Inches(3), width=Inches(8), height=Inches(1))
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        add_markdown_text(p, "Thank You!")
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(40)
        p.font.color.rgb = RGBColor(48, 40, 36)

        # Regex to remove **...**
        bold_pattern = re.compile(r"\*\*")


        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = bold_pattern.sub("", run.text)


        # Save the presentation
        pptx_path = os.path.join(self.output_dir, filename)
        pres.save(pptx_path)
        print(f"\n✅ PowerPoint slides saved to {pptx_path}")
        return pptx_path



# ----------------------------
# Example usage
# ----------------------------
if __name__ == "__main__":
    model_output = """Slide 1: Title Slide
- Title: AI in Industry: Transforming Business Operations
- Subtitle: Awareness Presentation for Business Executives
- Author: Amey G M
- Date: October 26, 2023
- Tagline: Unlocking Efficiency and Innovation

Slide 2: AI's Transformative Impact
- Automating Tasks: AI streamlines operations by automating repetitive processes.
- Improving Decision-Making: AI-powered analytics provide data-driven insights.
- Enabling New Products & Services: AI facilitates innovative products and services.

"""

    generator = SlideGeneratorGold()
    slides_data = generator.parse_model_output(model_output)
    generator.create_slide_deck(slides_data)

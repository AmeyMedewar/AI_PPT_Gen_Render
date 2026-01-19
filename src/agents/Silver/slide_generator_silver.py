import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import re

# ------------------------------
# Slide Generator with Formatting
# ------------------------------
class SlideGeneratorSilver:
    def __init__(self, output="Output/"):
        self.output_dir = output
        os.makedirs(self.output_dir, exist_ok=True)

    # ------------------------------
    # Parser for model output
    # ------------------------------
    def parse_model_output(self, text):
        """
        Convert LLM-generated text into a list of slides.
        Each slide is a dict: {"title": ..., "subtitle": ..., "tagline": ..., "bullets": [...]}
        """
        slides = []
        current_slide = None

        for line in text.splitlines():
            line = line.strip()
            if not line:
                continue

            # New slide
            if line.lower().startswith("slide") and ":" in line:
                if current_slide:
                    slides.append(current_slide)
                current_slide = {"title": "", "subtitle": "", "tagline": "", "bullets": []}
                slide_title = line.split(":", 1)[1].strip()
                current_slide["title"] = slide_title

            # Title / Subtitle / Tagline (overwrite if provided)
            elif line.lower().startswith("- title:"):
                current_slide["title"] = line.split(":",1)[1].strip()
            elif line.lower().startswith("- subtitle:"):
                current_slide["subtitle"] = line.split(":",1)[1].strip()
            elif line.lower().startswith("- tagline:"):
                current_slide["tagline"] = line.split(":",1)[1].strip()
            elif line.startswith("- "):
                # Skip title/subtitle/tagline lines
                if any(line.lower().startswith(f"- {k}:") for k in ["title","subtitle","tagline"]):
                    continue
                current_slide["bullets"].append(line[2:].strip())

        if current_slide:
            slides.append(current_slide)

        return slides

    # ------------------------------
    # Slide deck creation
    # ------------------------------
    def create_slide_deck(self, slides_data, filename="generated_ppt.pptx"):
        pres = Presentation()
        filename = slides_data[0].get("title"+".pptx","generated_ppt.pptx")

        for idx, slide_info in enumerate(slides_data):
            layout = pres.slide_layouts[0] if idx == 0 else pres.slide_layouts[1]
            slide = pres.slides.add_slide(layout)

            # ----------------------
            # Title formatting
            # ----------------------
            title_shape = slide.shapes.title
            title_shape.text = slide_info.get("title", "Untitled Slide")
            title_para = title_shape.text_frame.paragraphs[0]
            title_para.font.size = Pt(40)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(0, 51, 102)  # dark blue

            # ----------------------
            # Subtitle / Tagline (only first slide)
            # ----------------------
            if idx == 0:
                subtitle = slide_info.get("subtitle", "")
                tagline = slide_info.get("tagline", "")
                if subtitle or tagline:
                    placeholder = slide.placeholders[1]
                    placeholder.text = f"{subtitle}\n{tagline}".strip()
                    for p in placeholder.text_frame.paragraphs:
                        p.font.size = Pt(24)
                        p.font.italic = True
                        p.font.color.rgb = RGBColor(51, 51, 51)  # dark gray

            # ----------------------
            # Content bullets
            # ----------------------
            if idx != 0:
                content_placeholder = slide.placeholders[1]
                tf = content_placeholder.text_frame
                tf.clear()
                for bullet in slide_info.get("bullets", []):
                    for line in self._split_bullet(bullet):
                        p = tf.add_paragraph()
                        p.text = line
                        p.font.size = Pt(20)
                        p.font.bold = False
                        p.font.color.rgb = RGBColor(0, 0, 0)
                        p.level = 0  # can increase for sub-bullets

        # Regex to remove **...**
        bold_pattern = re.compile(r"\*\*")


        for slide in pres.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            run.text = bold_pattern.sub("", run.text)


        # Save the presentation
        pptx_path = os.path.join(self.output_dir, filename)
        pres.save(pptx_path)
        print(f"\n✅ PowerPoint slides saved to {pptx_path}")
        return pptx_path

    # ----------------------
    # Helper: Split long bullets
    # ----------------------
    def _split_bullet(self, text, max_chars=80):
        if len(text) <= max_chars:
            return [text]
        lines = []
        words = text.split()
        current = ""
        for w in words:
            if len(current) + len(w) + 1 <= max_chars:
                current += " " + w if current else w
            else:
                lines.append(current)
                current = w
        if current:
            lines.append(current)
        return lines


# ------------------------------
# Example Usage
# ------------------------------
if __name__ == "__main__":
    model_output = """Slide 1: Title Slide
- Title: AI in Industry
- Subtitle: Awareness Presentation
- Tagline: Transforming the Future of Work with Intelligence

Slide 2: Introduction
- AI is reshaping industries worldwide
- Applications in automation, analytics, and decision-making
- Huge potential for cost savings and efficiency

Slide 3: Key Applications
- Medical imaging (X-rays, MRI analysis)
- Predictive analytics for patient outcomes
- Personalized treatment plans
"""

    generator = SlideGeneratorSilver()
    slides_data = generator.parse_model_output(model_output)
    generator.create_slide_deck(slides_data)

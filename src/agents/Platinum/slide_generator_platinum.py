import os
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
import regex as re


class SlideGeneratorPlatinum:
    def __init__(self, output="Output/"):
        self.output_dir = output
        os.makedirs(self.output_dir, exist_ok=True)

    def parse_model_output(self, text):
        slides = []
        current_slide = None

        for line in text.splitlines():
            line = line.strip()
            if not line:
                continue

            if line.lower().startswith("slide") and ":" in line:
                if current_slide:
                    slides.append(current_slide)
                current_slide = {
                    "type": "normal",  # default
                    "title": "", "subtitle": "", "author": "",
                    "date": "", "tagline": "", "bullets": []
                }
                slide_title = line.split(":", 1)[1].strip()
                current_slide["title"] = slide_title

                # Detect special slide types
                lowered_title = slide_title.lower()
                if "summary" in lowered_title:
                    current_slide["type"] = "summary"
                elif "timeline" in lowered_title:
                    current_slide["type"] = "timeline"
                elif "flow" in lowered_title:
                    current_slide["type"] = "flow"
                elif "question" in lowered_title or "questions" in lowered_title:
                    current_slide["type"] = "question"
                elif "vs" in lowered_title or "comparison" in lowered_title or "difference" in lowered_title:
                    current_slide["type"] = "comparison"
                elif "thank you" in lowered_title:
                    current_slide["type"] = "thankyou"
                continue

            elif line.lower().startswith("- title:"):
                current_slide["title"] = line.split(":", 1)[1].strip()
            elif line.lower().startswith("- subtitle:"):
                current_slide["subtitle"] = line.split(":", 1)[1].strip()
            elif line.lower().startswith("- author:"):
                current_slide["author"] = line.split(":", 1)[1].strip()
            elif line.lower().startswith("- date:"):
                current_slide["date"] = line.split(":", 1)[1].strip()
            elif line.lower().startswith("- tagline:"):
                current_slide["tagline"] = line.split(":", 1)[1].strip()
            elif line.startswith("- "):
                bullet = line[2:].strip()
                current_slide["bullets"].append(bullet)

        if current_slide:
            slides.append(current_slide)

        return slides

    def select_image(self,title: str) -> str:
        """Return image path based on keywords in title."""
        title_lower = title.lower()

        if "cricket" in title_lower:
            return r"C:\Users\Amey\OneDrive\Desktop\Agentic_AI_Project\Data\default.jpg"
        elif "ai" in title_lower or "artificial intelligence" in title_lower:
            return r"C:\Users\Amey\OneDrive\Desktop\AI_based_ppt_genrator\AI.jpg"
        else:
            return r"C:\Users\Amey\OneDrive\Desktop\AI_based_ppt_genrator\Data\default.jpg"

    def add_title_slide(self,pres, slide_info):
        slide = pres.slides.add_slide(pres.slide_layouts[6])  # Blank layout

        # --- Title (Top Left, fixed box) ---
        title_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(0.7),
            width=Inches(7.5),   # fixed 60% width
            height=Inches(1.8),  # fixed height
        )
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        p_title = tf_title.paragraphs[0]
        p_title.text = slide_info.get("title", "Untitled Slide")
        p_title.font.size = Pt(44)  # starting size
        p_title.font.bold = True
        p_title.font.color.rgb = RGBColor(0, 51, 102)
        p_title.alignment = PP_ALIGN.LEFT

        # --- Subtitle (Below Title, fixed box) ---
        subtitle = slide_info.get("subtitle", "")
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                left=Inches(0.5),
                top=Inches(2.7),
                width=Inches(7.5),
                height=Inches(1.5),
            )
            tf_sub = subtitle_box.text_frame
            tf_sub.word_wrap = True
            tf_sub.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            p_sub = tf_sub.paragraphs[0]
            p_sub.text = subtitle
            p_sub.font.size = Pt(28)
            p_sub.font.color.rgb = RGBColor(47, 79, 79)
            p_sub.alignment = PP_ALIGN.LEFT

        # --- Author & Date (Bottom Left, fixed box) ---
        author = slide_info.get("author", "Amey G M")
        date = slide_info.get("date", "September 4, 2025")
        author_text = f"Author: {author}\nDate: {date}"

        meta_box = slide.shapes.add_textbox(
            left=Inches(0.5),
            top=Inches(6.0),
            width=Inches(4.5),
            height=Inches(1.0),
        )
        tf_meta = meta_box.text_frame
        tf_meta.word_wrap = True
        tf_meta.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf_meta.clear()
        for line in author_text.split("\n"):
            p = tf_meta.add_paragraph()
            p.text = line
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(59, 59, 59)
            p.alignment = PP_ALIGN.LEFT

        # --- Tagline (Bottom Center, fixed box) ---
        tagline = slide_info.get("tagline", "")
        if tagline:
            tagline_box = slide.shapes.add_textbox(
                left=Inches(0.5),
                top=Inches(7.0),
                width=Inches(7.5),
                height=Inches(0.6),
            )
            tf_tag = tagline_box.text_frame
            tf_tag.word_wrap = True
            tf_tag.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

            p_tag = tf_tag.paragraphs[0]
            p_tag.text = tagline
            p_tag.font.size = Pt(18)
            p_tag.font.color.rgb = RGBColor(0, 51, 102)
            p_tag.alignment = PP_ALIGN.CENTER

        # --- Image (Right 40%, full height, fixed) ---
        image_path = self.select_image(slide_info.get("title", "Untitled Slide"))
        if image_path:
            slide.shapes.add_picture(
                image_path,
                left=Inches(8.0),   # start after 60% text area (8 inches)
                top=0,
                width=Inches(5.3), # ~40% width of standard slide
                height=Inches(7.5) # full slide height
            )


    def add_normal_slide(self, pres, slide_info):
        slide = pres.slides.add_slide(pres.slide_layouts[1])

        # Set the slide title
        title_shape = slide.shapes.title
        title_shape.text = slide_info.get("title", "Untitled Slide")

        # Find the content placeholder (usually idx != 0)
        content_placeholder = None
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx != 0 and hasattr(placeholder, "text_frame"):
                content_placeholder = placeholder
                break

        def add_bold_text_paragraph(tf, text, font_size, level):
            """
            Add a paragraph to tf where text may contain **bold** parts.
            """
            # Split text by ** to detect bold segments
            parts = re.split(r'(\*\*)', text)
            p = tf.add_paragraph()
            p.level = level
            p.font.size = Pt(font_size)

            bold_mode = False
            for part in parts:
                if part == "**":
                    # Toggle bold mode on/off
                    bold_mode = not bold_mode
                elif part:
                    run = p.add_run()
                    run.text = part
                    run.font.size = Pt(font_size)
                    run.font.bold = bold_mode

        if content_placeholder:
            tf = content_placeholder.text_frame
            tf.clear()

            for bullet in slide_info.get("bullets", []):
                add_bold_text_paragraph(tf, bullet.strip(), 20, 0)


    def add_summary_slide(self, pres, slide_info):
        slide = pres.slides.add_slide(pres.slide_layouts[1])

        # Title of the slide
        slide.shapes.title.text = slide_info.get("title", "Summary")

        # Get the content placeholder
        tf = slide.placeholders[1].text_frame
        tf.clear()

        # Summary header
        p_intro = tf.add_paragraph()
        p_intro.text = "Summary Points:"
        p_intro.font.size = Pt(24)
        p_intro.font.bold = True
        p_intro.font.color.rgb = RGBColor(0, 51, 102)
        p_intro.space_after = Pt(10)
        p_intro.level = 0

        # Regex to remove the "**Key Point X:**" prefix
        prefix_pattern = re.compile(r"\*\*Key Point \d+:\*\*\s*")
        # Regex to strip all remaining ** from text
        bold_pattern = re.compile(r"\*\*(.*?)\*\*")

        # Actual bullet points
        for point in slide_info.get("bullets", []):
            # Remove the prefix
            clean_point = prefix_pattern.sub("", point).strip()
            # Remove any remaining ** markers
            clean_point = bold_pattern.sub(r"\1", clean_point)

            para = tf.add_paragraph()
            para.text = clean_point
            para.font.size = Pt(20)
            para.level = 0
            para.space_after = Pt(5)

    def add_timeline_slide(self, pres, slide_info):
        slide = pres.slides.add_slide(pres.slide_layouts[1])
        slide.shapes.title.text = slide_info.get("title", "Timeline")

        tf = slide.placeholders[1].text_frame
        tf.clear()

        # Regex to strip **...** markers
        bold_pattern = re.compile(r"\*\*(.*?)\*\*")

        for event in slide_info.get("bullets", []):
            # Remove any **...** while keeping text
            event = bold_pattern.sub(r"\1", event)

            if ":" in event:
                year, detail = event.split(":", 1)

                para = tf.add_paragraph()
                para.text = f"{year.strip()}:"
                para.font.size = Pt(20)
                para.font.bold = True
                para.font.color.rgb = RGBColor(0, 51, 102)
                para.level = 0

                detail_para = tf.add_paragraph()
                detail_para.text = detail.strip()
                detail_para.font.size = Pt(18)
                detail_para.level = 1
            else:
                para = tf.add_paragraph()
                para.text = event.strip()
                para.font.size = Pt(20)
                para.level = 0

    def add_flow_slide(self, pres, slide_info):
        slide = pres.slides.add_slide(pres.slide_layouts[6])  # Use blank layout

        # Add Title at the top
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = slide_info.get("title", "Flow")
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(36)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(0, 51, 102)

        # Add flow steps without bullet points
        flow_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        tf = flow_box.text_frame
        tf.word_wrap = True
        tf.clear()

        steps = slide_info.get("bullets", [])
        for i, step in enumerate(steps):
            symbol = "→" if i < len(steps) - 1 else "⏹"  # Arrow for steps, stop for last
            para = tf.add_paragraph()
            para.text = f"{symbol} {step.strip()}"
            para.font.size = Pt(20)
            para.font.color.rgb = RGBColor(0, 51, 102)
            para.level = 0
            para.space_after = Pt(8)

    def add_question_slide(self, pres, slide_info):
        slide = pres.slides.add_slide(pres.slide_layouts[6])  # Blank layout

        # Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = slide_info.get("title", "Questions")
        title_frame.paragraphs[0].font.size = Pt(36)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

        # Questions block (no bullets)
        questions_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        tf = questions_box.text_frame
        tf.word_wrap = True
        tf.clear()

        for idx, question in enumerate(slide_info.get("bullets", [])):
            if idx == 0:
                tf.text = f"❓ {question.strip()}"
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
                p.text = f"❓ {question.strip()}"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(0, 51, 102)
            p.level = 0
            p.space_after = Pt(10)

    
        
    def add_comparison_slide(self, pres, slide_info):
        slide = pres.slides.add_slide(pres.slide_layouts[1])  # Title and Content layout
        slide.shapes.title.text = slide_info.get("title", "Comparison")

        tf = slide.placeholders[1].text_frame
        tf.clear()

        def add_formatted_paragraph(text, level=0, font_size=20, bold=False):
            p = tf.add_paragraph()
            p.level = level
            p.font.size = Pt(font_size)
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size)
            run.font.bold = bold

        bullets = slide_info.get("bullets", [])

        # Regex to strip **...** markers
        bold_pattern = re.compile(r"\*\*(.*?)\*\*")

        for bullet in bullets:
            # Remove the "Aspect X:" prefix
            bullet = re.sub(r'^Aspect \d+:\s*', '', bullet)

            # Remove any **...** while keeping text
            bullet = bold_pattern.sub(r"\1", bullet)

            # If colon present, split term and explanation
            if ":" in bullet:
                term, explanation = bullet.split(":", 1)
                add_formatted_paragraph(term.strip() + ":", level=0, font_size=20, bold=True)
                add_formatted_paragraph(explanation.strip(), level=1, font_size=18, bold=False)
            else:
                add_formatted_paragraph(bullet.strip(), level=0, font_size=20, bold=True)
                            
    def add_thankyou_slide(self, pres, slide_info=None):
        slide = pres.slides.add_slide(pres.slide_layouts[6])
        
        # Big "Thank You"
        textbox = slide.shapes.add_textbox(
            left=Inches(1), 
            top=Inches(2), 
            width=Inches(8), 
            height=Inches(1.5)
        )
        tf = textbox.text_frame
        p = tf.paragraphs[0]
        p.text = "Thank You"
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)

        # Thank-you related tagline
        tagline_box = slide.shapes.add_textbox(
            left=Inches(1), 
            top=Inches(3.5),   # below the main text
            width=Inches(8), 
            height=Inches(1)
        )
        tagline_tf = tagline_box.text_frame
        p_tagline = tagline_tf.paragraphs[0]
        p_tagline.text = "We appreciate your time and attention"
        p_tagline.alignment = PP_ALIGN.CENTER
        p_tagline.font.size = Pt(20)
        p_tagline.font.italic = True
        p_tagline.font.color.rgb = RGBColor(80, 80, 80)  # soft gray


    def create_slide_deck(self, slides_data, filename="generated_ppt.pptx"):
        pres = Presentation()
        filename = slides_data[0].get("title","generated_ppt")+".pptx"


        for idx, slide_info in enumerate(slides_data):
            slide_type = slide_info.get("type", "normal").lower()

            if idx == 0:
                self.add_title_slide(pres, slide_info)
            elif slide_type == "normal":
                self.add_normal_slide(pres, slide_info)
            elif slide_type == "summary":
                self.add_summary_slide(pres, slide_info)
            elif slide_type == "timeline":
                self.add_timeline_slide(pres, slide_info)
            elif slide_type == "flow":
                self.add_flow_slide(pres, slide_info)
            elif slide_type == "question":
                self.add_question_slide(pres, slide_info)
            elif slide_type == "comparison":
                self.add_comparison_slide(pres, slide_info)
            elif slide_type == "thankyou":
                self.add_thankyou_slide(pres)
            else:
                print(f"⚠️ Unknown slide type: {slide_type}. Falling back to normal.")
                self.add_normal_slide(pres, slide_info)

        # Add final "Thank You" slide if not already included
        if not any(s.get("type") == "thankyou" for s in slides_data):
            self.add_thankyou_slide(pres)

        pptx_path = os.path.join(self.output_dir, filename)
        pres.save(pptx_path)
        print(f"\n✅ PowerPoint slides saved to {pptx_path}")
        return pptx_path


# ----------------------------
# Example usage
# ----------------------------
if __name__ == "__main__":

    model_output = """ Slide 1: Title Slide

- Title: ML vs DL
- Subtitle: Awareness Presentation for Business Executives
- Author: Amey G M
- Date: October 26, 2023
- Tagline:  Unlocking the Potential of Intelligent Systems


Slide 2: Understanding Machine Learning (ML)

- **Machine Learning (ML):**  A broad field encompassing algorithms that allow computers to learn from data without explicit programming.
- **Focus:** Identifying patterns, making predictions, and improving performance over time based on data analysis.
- **Examples:** Spam filtering, recommendation systems, fraud detection.
- **Key Advantage:**  Relatively simpler to implement and requires less data compared to Deep Learning in some cases.
- **Limitations:**  May struggle with complex, unstructured data; requires feature engineering.


Slide 3: Deep Dive into Deep Learning (DL)

- **Deep Learning (DL):** A subfield of ML utilizing artificial neural networks with multiple layers (hence "deep").
- **Focus:**  Learning complex patterns and representations from raw data, often requiring vast datasets.
- **Examples:** Image recognition, natural language processing, self-driving cars.
- **Key Advantage:**  Exceptional performance on complex tasks; automates feature extraction.
- **Limitations:**  Requires significant computational resources and large datasets; can be a "black box" in terms of interpretability.


Slide 4: ML vs. DL: A Comparison

- **Aspect 1: Data Requirements:** ML - _Moderate_, DL - _Large_
- **Aspect 2: Computational Power:** ML - _Lower_, DL - _High_
- **Aspect 3: Complexity:** ML - _Simpler_, DL - _More Complex_
- **Aspect 4: Interpretability:** ML - _More Interpretable_, DL - _Less Interpretable_
- **Aspect 5: Application Suitability:** ML - Suitable for simpler tasks with structured data; DL - Suitable for complex tasks with unstructured data.


Slide 5: Summary / Key Takeaways

- **Key Point 1:** Machine Learning (ML) is a broader field encompassing algorithms that enable computers to learn from data.      
- **Key Point 2:** Deep Learning (DL) is a specialized subset of ML using deep neural networks for complex pattern recognition.    
- **Key Point 3:** DL offers superior performance on complex tasks but demands substantial computational resources and data.       
- **Key Point 4:** ML is often more interpretable and requires less data, making it suitable for simpler applications.
- **Key Point 5:** The choice between ML and DL depends on the specific business problem, available data, and computational capabilities.


Slide 6: Thank You

- Thank you for your attention.  Please contact me with any questions.  [Amey G M Contact Information]

"""

    generator = SlideGeneratorPlatinum()
    slides_data = generator.parse_model_output(model_output)
    print(slides_data)
    generator.create_slide_deck(slides_data)
